import os
import json
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext in [".txt", ".md", ".log"]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    if ext == ".json":
        with open(file_path, "r", encoding="utf-8") as f:
            return json.dumps(json.load(f), indent=2)

    if ext == ".csv":
        return pd.read_csv(file_path).to_string()

    if ext in [".xls", ".xlsx"]:
        return pd.read_excel(file_path).to_string()

    if ext == ".pdf":
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text)
        return "\n".join(text)

    if ext == ".docx":
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    raise ValueError(f"Unsupported file type: {ext}")


def chunk_text(text, chunk_size=300, overlap=50):
    words = text.split()

    if not words:
        return []

    chunks = []
    start = 0

    while start < len(words):
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
        start += chunk_size - overlap

    return chunks